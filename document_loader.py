import os
import re
from typing import List, Dict, Any
from pypdf import PdfReader
from google import genai

# Extension categories
TEXT_EXTENSIONS = ['.txt', '.md', '.java', '.py', '.js', '.html', '.css', '.c', '.cpp', '.h', '.json', '.xml', '.csv']
PDF_EXTENSIONS = ['.pdf']
IMAGE_EXTENSIONS = ['.png', '.jpg', '.jpeg', '.webp', '.bmp']
EXCEL_EXTENSIONS = ['.xlsx', '.xls']
PPT_EXTENSIONS = ['.pptx']

def clean_text(text: str) -> str:
    """Removes extra whitespaces and cleans up text."""
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def load_text_file(file_path: str) -> str:
    """Loads content from a plain text or source code file."""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def load_pdf_file(file_path: str) -> str:
    """Loads content from a PDF file."""
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

def extract_text_via_gemini(file_path: str) -> str:
    """Uses Gemini 2.5 Flash to extract text (OCR) from scanned documents/images."""
    import shutil
    temp_path = None
    try:
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key or api_key == "your_gemini_api_key_here":
            print(f"[OCR WARNING] '{os.path.basename(file_path)}' 파일 처리를 위한 GEMINI_API_KEY가 설정되지 않았습니다. (.env 파일을 확인해 주세요)")
            return ""
            
        client = genai.Client(api_key=api_key)
        
        # Workaround for google-genai SDK unicode header bug (httpx encodes headers in ASCII)
        filename = os.path.basename(file_path)
        has_non_ascii = any(ord(char) > 127 for char in filename)
        
        upload_path = file_path
        if has_non_ascii:
            dir_name = os.path.dirname(file_path)
            ext = os.path.splitext(file_path)[1]
            temp_path = os.path.join(dir_name, f"temp_ocr_upload{ext}")
            shutil.copy2(file_path, temp_path)
            upload_path = temp_path
            
        # Upload the file to Gemini File API
        print(f"[OCR] Gemini API를 사용하여 '{filename}' 파일의 이미지 텍스트를 추출하는 중...")
        uploaded_file = client.files.upload(file=upload_path)
        
        response = client.models.generate_content(
            model="gemini-flash-latest",
            contents=[
                uploaded_file,
                "이 이미지/문서 파일에 적힌 모든 텍스트를 누락 없이 한글로 정확하게 읽어서 그대로 텍스트로 추출해 주세요. 부연 설명이나 메타 정보는 제외하고 오직 추출된 본문 텍스트만 출력해 주세요."
            ]
        )
        
        # Clean up the file from Gemini cloud storage
        try:
            client.files.delete(name=uploaded_file.name)
        except Exception:
            pass
            
        extracted = response.text.strip()
        print(f"[OCR SUCCESS] Gemini OCR로 '{filename}'에서 {len(extracted)}자 추출 완료.")
        return extracted
    except Exception as e:
        print(f"[OCR ERROR] Gemini OCR 처리 중 오류 발생: {e}")
        return ""
    finally:
        # Clean up the temporary local file if created
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except Exception:
                pass

def load_excel_file(file_path: str) -> str:
    """Loads content from an Excel file, converting sheets to text CSV tables."""
    try:
        import pandas as pd
        xls = pd.ExcelFile(file_path)
        text_parts = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            df = df.dropna(how='all')
            if df.empty:
                continue
            csv_str = df.to_csv(index=False)
            text_parts.append(f"[시트: {sheet_name}]\n{csv_str}")
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"Error loading Excel file {file_path}: {e}")
        return ""

def load_pptx_file(file_path: str) -> str:
    """Loads text content from a PowerPoint presentation file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"[슬라이드 {idx+1}]\n" + "\n".join(slide_text))
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"Error loading PowerPoint file {file_path}: {e}")
        return ""

def split_text(text: str, chunk_size: int = 500, chunk_overlap: int = 100) -> List[str]:
    """Splits text into chunks of specified size and overlap."""
    chunks = []
    start = 0
    text_len = len(text)
    
    if text_len <= chunk_size:
        return [text]
        
    while start < text_len:
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start += chunk_size - chunk_overlap
        
    return chunks

def load_and_chunk_documents(dir_path: str, chunk_size: int = 500, chunk_overlap: int = 100) -> List[Dict[str, Any]]:
    """Loads PDF, TXT, Excel, PPTX, Images, and Source files in a directory and chunks them."""
    documents = []
    
    if not os.path.exists(dir_path):
        os.makedirs(dir_path, exist_ok=True)
        return documents
        
    for root, _, files in os.walk(dir_path):
        for file in files:
            file_path = os.path.join(root, file)
            ext = os.path.splitext(file)[1].lower()
            
            try:
                content = ""
                if ext in TEXT_EXTENSIONS:
                    content = load_text_file(file_path)
                elif ext in PDF_EXTENSIONS:
                    content = load_pdf_file(file_path)
                    content_clean = clean_text(content)
                    # If empty or extremely short, fall back to Gemini OCR
                    if len(content_clean) < 10:
                        print(f"[INFO] '{file}' 파일에서 일반 텍스트를 추출할 수 없습니다. Gemini OCR을 시도합니다...")
                        content = extract_text_via_gemini(file_path)
                elif ext in IMAGE_EXTENSIONS:
                    print(f"[INFO] '{file}' 이미지 파일 감지. Gemini OCR을 시도합니다...")
                    content = extract_text_via_gemini(file_path)
                elif ext in EXCEL_EXTENSIONS:
                    content = load_excel_file(file_path)
                elif ext in PPT_EXTENSIONS:
                    content = load_pptx_file(file_path)
                else:
                    continue
                
                content = clean_text(content)
                if not content:
                    print(f"[WARNING] '{file}' 파일에서 텍스트를 추출하지 못했습니다. (텍스트 레이어가 없는 스캔된 이미지 PDF이거나 비어 있는 파일인지 확인해 주세요.)")
                    continue
                    
                chunks = split_text(content, chunk_size, chunk_overlap)
                
                for idx, chunk in enumerate(chunks):
                    documents.append({
                        "id": f"{file}_{idx}",
                        "content": chunk,
                        "metadata": {
                            "source": file,
                            "path": file_path,
                            "chunk_index": idx
                        }
                    })
                print(f"Loaded and chunked: {file} into {len(chunks)} chunks")
            except Exception as e:
                print(f"Error reading {file}: {e}")
                
    return documents

if __name__ == "__main__":
    # Test document loader
    test_dir = "./test_data"
    os.makedirs(test_dir, exist_ok=True)
    with open(os.path.join(test_dir, "test.txt"), "w", encoding="utf-8") as f:
        f.write("안녕하세요! 이것은 Gemini RAG 프로젝트의 테스트 문서입니다. " * 20)
        
    docs = load_and_chunk_documents(test_dir)
    for doc in docs:
        print(f"ID: {doc['id']}")
        print(f"Metadata: {doc['metadata']}")
        print(f"Content: {doc['content'][:100]}...\n")
